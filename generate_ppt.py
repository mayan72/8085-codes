#!/usr/bin/env python3
"""Build a polished 16:9 deck from existing 8085 program notes. Content is preserved."""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import nsmap, qn
from pptx.oxml import parse_xml
from pptx.util import Emu, Inches, Pt
from lxml import etree

# --- Theme ---
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
NAVY_MID = RGBColor(0x14, 0x32, 0x55)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
GOLD_SOFT = RGBColor(0xE8, 0xC9, 0x5A)
CREAM = RGBColor(0xF6, 0xF1, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x24, 0x33)
MUTED = RGBColor(0x5B, 0x67, 0x75)
CODE_BG = RGBColor(0x0D, 0x16, 0x24)
CODE_FG = RGBColor(0xE6, 0xED, 0xF3)
CODE_MUTED = RGBColor(0x8B, 0x9C, 0xB0)
TEAL = RGBColor(0x2B, 0xB3, 0xA8)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_run_font(run, name, size, color, bold=False):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", name)


def add_rect(slide, l, t, w, h, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    spPr = shape._element.spPr
    # no shadow
    return shape


def add_round_rect(slide, l, t, w, h, fill, radius=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    # adjust corner
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass
    return shape


def add_textbox(slide, l, t, w, h, text, font, size, color, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        box.text_frame._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}.get(anchor, "t"))
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, font, size, color, bold)
    return box


def add_multiline(slide, l, t, w, h, lines, font, size, color, bold=False, spacing=1.05, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(2)
        p.line_spacing = spacing
        run = p.add_run()
        run.text = line
        set_run_font(run, font, size, color, bold)
    return box


def footer(slide, page, total, dark=False):
    color = RGBColor(0x9A, 0xA8, 0xB8) if dark else MUTED
    add_textbox(slide, Inches(0.5), Inches(7.18), Inches(9), Inches(0.28),
                "8085 Microprocessor Programs  ·  Original code unchanged",
                "Calibri", 11, color, False)
    add_textbox(slide, Inches(11.4), Inches(7.18), Inches(1.4), Inches(0.28),
                f"{page}  /  {total}", "Calibri", 11, color, False, PP_ALIGN.RIGHT)


def chrome_content(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, CREAM)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), GOLD)
    add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, NAVY)
    add_rect(slide, 0, Inches(7.05), SLIDE_W, Inches(0.45), RGBColor(0xEE, 0xE8, 0xDC))
    add_textbox(slide, Inches(0.55), Inches(0.22), Inches(12.2), Inches(0.48),
                title, "Calibri", 26, NAVY, True)
    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(0.68), Inches(12.2), Inches(0.32),
                    subtitle, "Calibri", 13, MUTED, False)


def chrome_dark(slide):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.1), GOLD)


def code_card(slide, l, t, w, h, lines, title="Code"):
    card = add_round_rect(slide, l, t, w, h, CODE_BG, 0.05)
    add_rect(slide, l, t, w, Inches(0.38), RGBColor(0x12, 0x22, 0x36))
    for i, col in enumerate((RGBColor(0xFF, 0x5F, 0x56), RGBColor(0xFF, 0xBD, 0x2E), RGBColor(0x27, 0xC9, 0x3F))):
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, l + Inches(0.18 + i * 0.22), t + Inches(0.12), Inches(0.14), Inches(0.14))
        d.fill.solid()
        d.fill.fore_color.rgb = col
        d.line.fill.background()
    add_textbox(slide, l + Inches(0.95), t + Inches(0.04), w - Inches(1.15), Inches(0.3),
                title, "Calibri", 12, GOLD_SOFT, True)
    box = slide.shapes.add_textbox(l + Inches(0.22), t + Inches(0.46), w - Inches(0.4), h - Inches(0.58))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(1)
        p.line_spacing = 1.05
        run = p.add_run()
        run.text = line if line != "" else " "
        set_run_font(run, "Consolas", 13, CODE_FG, False)
    return card


def note_card(slide, l, t, w, h, title, lines):
    add_round_rect(slide, l, t, w, h, WHITE, 0.06)
    # left gold strip
    add_rect(slide, l, t, Inches(0.08), h, GOLD)
    add_textbox(slide, l + Inches(0.25), t + Inches(0.12), w - Inches(0.4), Inches(0.32),
                title, "Calibri", 14, NAVY, True)
    box = slide.shapes.add_textbox(l + Inches(0.25), t + Inches(0.46), w - Inches(0.42), h - Inches(0.58))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        p.line_spacing = 1.12
        run = p.add_run()
        run.text = line
        set_run_font(run, "Calibri", 13, INK, False)


def pill(slide, l, t, w, h, text):
    s = add_round_rect(slide, l, t, w, h, RGBColor(0x14, 0x32, 0x55), 0.5)
    add_textbox(slide, l, t, w, h, text, "Calibri", 11, GOLD_SOFT, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    return s


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # Pre-count slides for footer
    # We'll add then stamp later by iterating? Easier to hardcode total.
    # Structure:
    # 1 title, 2 agenda, then programs...
    slides_meta = []

    def new():
        s = prs.slides.add_slide(blank)
        slides_meta.append(s)
        return s

    # ===== 1 TITLE =====
    s = new()
    chrome_dark(s)
    add_rect(s, 0, Inches(5.9), SLIDE_W, Inches(1.6), NAVY_MID)
    add_rect(s, Inches(0.7), Inches(1.7), Inches(1.4), Inches(0.08), GOLD)
    add_textbox(s, Inches(0.7), Inches(1.9), Inches(11.8), Inches(1.1),
                "Assembly Language Programs", "Calibri", 40, WHITE, True)
    add_textbox(s, Inches(0.7), Inches(2.95), Inches(11.8), Inches(0.7),
                "8085 Microprocessor", "Calibri", 32, GOLD_SOFT, False)
    add_textbox(s, Inches(0.7), Inches(3.8), Inches(11), Inches(0.5),
                "A visual restyle of the original notes. Program text, comments, and explanations are unchanged.",
                "Calibri", 16, RGBColor(0xC5, 0xD0, 0xDC), False)
    add_textbox(s, Inches(0.7), Inches(6.2), Inches(11), Inches(0.35),
                "Codes  ·  Complements  ·  Arithmetic  ·  Sign & Parity checks",
                "Calibri", 16, WHITE, False)

    # ===== 2 AGENDA =====
    s = new()
    chrome_content(s, "Contents", "Programs included in this deck")
    items = [
        ("01", "1's and 2's complement", "8-bit and 16-bit numbers"),
        ("02", "Square of an 8-bit number", "Repeated addition using H-L"),
        ("03", "Factorial of a number", "Subroutine multiply by repeated add"),
        ("04", "Multiply two 8-bit numbers", "Using logical rotate instructions"),
        ("05", "Add two 16-bit numbers", "Low byte ADD, high byte ADC"),
        ("06", "Positive or negative number", "Sign bit via RAL and carry"),
        ("07", "Odd or even number", "LSB via RAR and carry"),
        ("08", "Subtract two 8-bit numbers", "With or without borrow"),
    ]
    for i, (num, title, sub) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.5 + col * 6.35)
        y = Inches(1.15 + row * 1.35)
        add_round_rect(s, x, y, Inches(6.05), Inches(1.18), WHITE, 0.08)
        add_rect(s, x, y, Inches(0.1), Inches(1.18), GOLD if i % 2 == 0 else TEAL)
        add_textbox(s, x + Inches(0.3), y + Inches(0.18), Inches(1.0), Inches(0.4),
                    num, "Calibri", 22, GOLD, True)
        add_textbox(s, x + Inches(1.25), y + Inches(0.18), Inches(4.5), Inches(0.4),
                    title, "Calibri", 18, NAVY, True)
        add_textbox(s, x + Inches(1.25), y + Inches(0.6), Inches(4.5), Inches(0.35),
                    sub, "Calibri", 13, MUTED, False)

    # ===== COMPLEMENT SECTION =====
    s = new()
    chrome_dark(s)
    add_textbox(s, Inches(0.7), Inches(2.4), Inches(3), Inches(0.4),
                "PROGRAM  01", "Calibri", 16, GOLD, True)
    add_textbox(s, Inches(0.7), Inches(2.9), Inches(12), Inches(1.0),
                "1's , 2's complement of 8 bit number", "Calibri", 34, WHITE, True)
    add_textbox(s, Inches(0.7), Inches(4.1), Inches(11), Inches(0.5),
                "Followed by 1's and 2's complement of 16 bit number", "Calibri", 18, RGBColor(0xC5, 0xD0, 0xDC), False)

    s = new()
    chrome_content(s, "1's , 2's complement of 8 bit number", "Code")
    code_card(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(5.55), [
        "1'S COMPLEMENT",
        "",
        "LDA C050H",
        "CMA",
        "STA C051H",
        "HLT",
        "",
        "2'S COMPLEMENT",
        "",
        "LDA C050H",
        "CMA",
        "INR A",
        "STA C051H",
        "HLT",
    ], "8-bit  ·  original listing")

    s = new()
    chrome_content(s, "1's , 2's complement of 16 bit number", "1'S COMPLEMENT")
    code_card(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(5.55), [
        "1'S COMPLEMENT",
        "",
        "LXI H C050H",
        "MOV A M",
        "CMA",
        "STA CO52H",
        "INX H",
        "MOV A M",
        "STA C053H",
        "HLT",
    ], "16-bit  ·  1'S COMPLEMENT  ·  original listing")

    s = new()
    chrome_content(s, "1's , 2's complement of 16 bit number", "2'COMPLEMENT")
    code_card(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(5.55), [
        "2'COMPLEMENT",
        "",
        "LXI H C050H",
        "MVI B 00H",
        "MOV A M",
        "CMA",
        "ADI 01H",
        "STA C052H",
        "JNC LOOP",
        "INR B",
        "LOOP LXI H C051H",
        "MOV A M",
        "CMA",
        "ADD B",
        "STA C053H",
        "HLT",
    ], "16-bit  ·  2'COMPLEMENT  ·  original listing")

    # ===== SQUARE =====
    s = new()
    chrome_dark(s)
    add_textbox(s, Inches(0.7), Inches(2.4), Inches(3), Inches(0.4),
                "PROGRAM  02", "Calibri", 16, GOLD, True)
    add_textbox(s, Inches(0.7), Inches(2.9), Inches(12), Inches(1.2),
                "8085 program to find square of a 8 bit number", "Calibri", 32, WHITE, True)

    s = new()
    chrome_content(s, "8085 program to find square of a 8 bit number", "code")
    code_card(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(5.55), [
        "2000    MVI H 20      H <- 20",
        "2002    MVI L 50      L <- 50",
        "2004    MVI A 00      A <- 00",
        "2006    MOV B, M      B <- M",
        "2007    ADD M         A <- A + M",
        "2008    DCR B         B <- B – 01",
        "2009    JNZ 2007      Jump if ZF = 0",
        "200C    STA 3050      M[3050] <- A",
        "200F    HLT           END",
    ], "original listing with addresses")

    s = new()
    chrome_content(s, "8085 program to find square of a 8 bit number", "Explanation – Registers used A, H, L, B and indirect memory M:")
    note_card(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(5.55), "Explanation", [
        "MVI H 20 – initialize register H with 20",
        "MVI L 50 – initialize register L with 50",
        "MVI A 00 – initialize accumulator A with 00",
        "MOV B, M – moves the content of memory location which is indirectly specified by M in register B",
        "ADD M – add the content of memory location which is indirectly specified by M in accumulator A",
        "DCR B – decrement value of register B by 1",
        "JNZ 2007 – jump to memory location 2007 if ZF = 0, i.e register B does not contain 0",
        "STA 3050 – stores value of A in 3050",
        "HLT – stops executing the program and halts any further execution",
    ])

    # ===== FACTORIAL =====
    s = new()
    chrome_dark(s)
    add_textbox(s, Inches(0.7), Inches(2.4), Inches(3), Inches(0.4),
                "PROGRAM  03", "Calibri", 16, GOLD, True)
    add_textbox(s, Inches(0.7), Inches(2.9), Inches(12), Inches(1.2),
                "8085 program to find the factorial of a number", "Calibri", 32, WHITE, True)

    s = new()
    chrome_content(s, "8085 program to find the factorial of a number", "code")
    code_card(s, Inches(0.4), Inches(1.1), Inches(12.5), Inches(5.65), [
        "2000H    Data                      Data Byte",
        "2001H    Result                    Result of factorial",
        "2002H    LXI H, 2000H              Load data from memory",
        "2005H    MOV B, M                  Load data to B register",
        "2006H    MVI D, 01H                Set D register with 1",
        "2008H    FACTORIAL  CALL MULTIPLY  Subroutine call for multiplication",
        "200BH    DCR B                     Decrement B",
        "200CH    JNZ FACTORIAL             Call factorial till B becomes 0",
        "200FH    INX H                     Increment memory",
        "2010H    MOV M, D                  Store result in memory",
        "2011H    HLT                       Halt",
        "2100H    MULTIPLY  MOV E, B        Transfer contents of B to C",
        "2101H    MVI A, 00H                Clear accumulator to store result",
        "2103H    MULTIPLYLOOP  ADD D       Add contents of D to A",
        "2104H    DCR E                     Decrement E",
        "2105H    JNZ MULTIPLYLOOP          Repeated addition",
        "2108H    MOV D, A                  Transfer contents of A to D",
        "2109H    RET                       Return from subroutine",
    ], "original listing")

    s = new()
    chrome_content(s, "8085 program to find the factorial of a number", "Explanation")
    note_card(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(5.55), "Explanation", [
        "First set register B with data.",
        "Set register D with data by calling MULTIPLY subroutine one time.",
        "Decrement B and add D to itself B times by calling MULTIPLY subroutine as 4*3 is equivalent to 4+4+4 (i.e., 3 times).",
        "Repeat the above step till B reaches 0 and then exit the program.",
        "The result is obtained in D register which is stored in memory",
    ])

    # ===== MULTIPLY LOGICAL =====
    s = new()
    chrome_dark(s)
    add_textbox(s, Inches(0.7), Inches(2.4), Inches(3), Inches(0.4),
                "PROGRAM  04", "Calibri", 16, GOLD, True)
    add_textbox(s, Inches(0.7), Inches(2.9), Inches(12), Inches(1.3),
                "8085 program to multiply two 8 bit numbers using logical instructions", "Calibri", 30, WHITE, True)

    s = new()
    chrome_content(s, "Multiply two 8 bit numbers using logical instructions", "code")
    code_card(s, Inches(0.5), Inches(1.15), Inches(6.15), Inches(5.55), [
        "MVI B 05      B <- 05",
        "MVI C 04      C <- 04",
        "MOV A, B      A <- B",
        "RLC           rotate the content of A without carry",
        "RLC           rotate the content of A without carry",
        "STA 3050      3050 <- A",
        "HLT           End of the program",
    ], "original listing")
    note_card(s, Inches(6.9), Inches(1.15), Inches(5.9), Inches(5.55), "Explanation", [
        "MVI B 05: assign the value 05 to B register.",
        "MVI C 04: assign the value 04 to C register.",
        "MOV A, B: move the content of register B to register A.",
        "RLC: rotate the content of accumulator left without carry.",
        "RLC: rotate the content of accumulator left without carry.",
        "STA 3050: store the content of register A at memory location 3050",
        "HLT: stops the execution of the program.",
    ])

    # ===== ADD 16 =====
    s = new()
    chrome_dark(s)
    add_textbox(s, Inches(0.7), Inches(2.4), Inches(3), Inches(0.4),
                "PROGRAM  05", "Calibri", 16, GOLD, True)
    add_textbox(s, Inches(0.7), Inches(2.9), Inches(12), Inches(1.2),
                "8085 program to add two 16 bit numbers", "Calibri", 32, WHITE, True)

    s = new()
    chrome_content(s, "8085 program to add two 16 bit numbers", "code")
    code_card(s, Inches(0.5), Inches(1.15), Inches(6.15), Inches(5.55), [
        "LDA 2050      A ← 2050",
        "MOV B, A      B ← A",
        "LDA 2052      A ← 2052",
        "ADD B         A ← A+B",
        "STA 3050      A → 3050",
        "LDA 2051      A ← 2051",
        "MOV B, A      B ← A",
        "LDA 2053      A ← 2053",
        "ADC B         A ← A+B+CY",
        "STA 3051      A → 3051",
        "HLT           Stops execution",
    ], "original listing")
    note_card(s, Inches(6.9), Inches(1.15), Inches(5.9), Inches(5.55), "Explanation:", [
        "LDA 2050 stores the value at 2050 in A (accumulator).",
        "MOV B, A stores the value of A into the B register.",
        "LDA 2052 stores the value at 2052 in A.",
        "ADD B add the contents of B and A and store them in A.",
        "STA 3050 stores the result in memory location 3050.",
        "LDA 2051 stores the value at 2051 in A.",
        "MOV B, A stores the value of A into the B register.",
        "LDA 2053 stores the value at 2053 in A.",
        "ADC B adds the contents of B, A, and carry from the lower bit addition and store in A.",
        "STA 3051 stores the result in memory location 3051.",
        "HLT stops execution.",
    ])

    # ===== POS/NEG =====
    s = new()
    chrome_dark(s)
    add_textbox(s, Inches(0.7), Inches(2.4), Inches(3), Inches(0.4),
                "PROGRAM  06", "Calibri", 16, GOLD, True)
    add_textbox(s, Inches(0.7), Inches(2.9), Inches(12), Inches(1.2),
                "Wearher a number is positive or negative", "Calibri", 32, WHITE, True)

    s = new()
    chrome_content(s, "Wearher a number is positive or negative", "Code")
    code_card(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(5.55), [
        "LXI H C050H",
        "MOV A M",
        "RAL",
        "JC LOOP",
        "MVI A 00H",
        "STA C053H",
        "HLT",
        "LOOP MVI A 01H",
        "STA C053H",
        "HLT",
    ], "original listing")

    # ===== ODD EVEN =====
    s = new()
    chrome_dark(s)
    add_textbox(s, Inches(0.7), Inches(2.4), Inches(3), Inches(0.4),
                "PROGRAM  07", "Calibri", 16, GOLD, True)
    add_textbox(s, Inches(0.7), Inches(2.9), Inches(12), Inches(1.2),
                "To find weather the number is odd or even", "Calibri", 32, WHITE, True)

    s = new()
    chrome_content(s, "To find weather the number is odd or even", "Code")
    code_card(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(5.55), [
        "LDA C050H",
        "RAR",
        "JC LOOP",
        "MVI A 00H",
        "STA C055H",
        "HLT",
        "LOOP MVI A 01H",
        "STA C055H",
        "HLT",
    ], "original listing")

    # ===== SUBTRACT =====
    s = new()
    chrome_dark(s)
    add_textbox(s, Inches(0.7), Inches(2.4), Inches(3), Inches(0.4),
                "PROGRAM  08", "Calibri", 16, GOLD, True)
    add_textbox(s, Inches(0.7), Inches(2.9), Inches(12), Inches(1.3),
                "8085 program to subtract two 8-bit numbers with or without borrow", "Calibri", 28, WHITE, True)

    s = new()
    chrome_content(s, "Subtract two 8-bit numbers with or without borrow", "code")
    code_card(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(5.55), [
        "MVI     C, 00      [C] <- 00",
        "LHLD    2500       [H-L] <- [2500]",
        "MOV     A, H       [A] <- [H]",
        "SUB     L          [A] <- [A] – [L]",
        "JNC     200B       Jump If no borrow",
        "INR     C          [C] <- [C] + 1",
        "STA     2502       [A] -> [2502], Result",
        "MOV     A, C       [A] <- [C]",
        "STA     2503       [A] -> [2503], Borrow",
        "HLT                Stop",
    ], "original listing")

    s = new()
    chrome_content(s, "Subtract two 8-bit numbers with or without borrow", "Explanation – Registers A, H, L, C are used for general purpose:")
    note_card(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(5.55), "Explanation", [
        "MOV is used to transfer the data from memory to accumulator (1 Byte)",
        "LHLD is used to load register pair directly using 16-bit address (3 Byte instruction)",
        "MVI is used to move data immediately into any of registers (2 Byte)",
        "STA is used to store the content of accumulator into memory(3 Byte instruction)",
        "INR is used to increase register by 1 (1 Byte instruction)",
        "JNC is used to jump if no borrow (3 Byte instruction)",
        "SUB is used to subtract two numbers where one number is in accumulator(1 Byte)",
        "HLT is used to halt the program",
    ])

    # ===== CLOSE =====
    s = new()
    chrome_dark(s)
    add_rect(s, Inches(0.7), Inches(2.55), Inches(1.4), Inches(0.08), GOLD)
    add_textbox(s, Inches(0.7), Inches(2.8), Inches(12), Inches(0.9),
                "End of programs", "Calibri", 36, WHITE, True)
    add_textbox(s, Inches(0.7), Inches(3.75), Inches(11.5), Inches(0.8),
                "All listings, comments, and explanations match the original text files in this repository.",
                "Calibri", 18, RGBColor(0xC5, 0xD0, 0xDC), False)

    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        # detect dark by looking at first rect fill? simpler: first slide and section dividers and last
        # We'll inspect background via first shape
        dark = False
        if slide.shapes:
            sh = slide.shapes[0]
            try:
                fc = sh.fill.fore_color.rgb
                dark = fc == NAVY
            except Exception:
                dark = False
        footer(slide, i, total, dark=dark)

    out = "/workspace/8085_Microprocessor_Programs.pptx"
    prs.save(out)
    print(f"Wrote {out} with {total} slides")


if __name__ == "__main__":
    build()
